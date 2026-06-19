"""Drone Safety System deck — warm paper / textbook theme.

  * beige-brown paper background (PIL)
  * elegant serif type (Constantia / Georgia)
  * light glassmorphism cards
  * hand-drawn, sketch-style diagrams (matplotlib sketch + Segoe Print)

Run:  ../../.venv/Scripts/python.exe build_deck.py
"""
from __future__ import annotations
import math
import os
import numpy as np
from PIL import Image
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch, Circle, Rectangle

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsdecls, qn
from pptx.oxml import parse_xml

HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(HERE, "assets"); os.makedirs(ASSETS, exist_ok=True)
SK = os.path.join(ASSETS, "sketches"); os.makedirs(SK, exist_ok=True)
PAPER_PNG = os.path.join(ASSETS, "paper.png")
DRONE_PNG = os.path.join(ASSETS, "drone_hero.png")

# ---- palette (light / minimal-mechanical, warm earth tones) ----
INK    = RGBColor(0x34, 0x30, 0x2A)   # warm dark brown (text)
SOFT   = RGBColor(0x5E, 0x57, 0x4C)   # secondary text
MUTED  = RGBColor(0x8C, 0x84, 0x74)   # tertiary / labels
TERRA  = RGBColor(0xB0, 0x74, 0x4D)   # clay accent
OLIVE  = RGBColor(0x7C, 0x8A, 0x5A)   # muted olive (success)
INDIGO = RGBColor(0x5B, 0x72, 0x86)   # muted slate-blue (primary)
MUST   = RGBColor(0xB5, 0x90, 0x4E)   # ochre (IP)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
PANEL  = RGBColor(0xFB, 0xF8, 0xF1)   # card fill (warm off-white)
HAIR   = RGBColor(0xCB, 0xC1, 0xAD)   # warm hairline border

F_HERO  = "Constantia"                # serif display (crafted feel)
F_TITLE = "Constantia"
F_BODY  = "Segoe UI"                   # clean sans body (editorial contrast)
F_MONO  = "Consolas"                  # technical labels / figure numbers

# diagram ink (hex, for matplotlib — ink-brown lines + muted accents on cream)
mINK = "#4A4338"; mTERRA = "#B0744D"; mINDIGO = "#5B7286"; mOLIVE = "#7C8A5A"
mMUST = "#B5904E"; mMUTE = "#968D7C"

SW, SH = Inches(13.333), Inches(7.5)


# ===================================================================
#  PAPER BACKGROUND
# ===================================================================
def gen_paper(path):
    """Warm cream sheet with a faint engineering grid, margin frame & corner
    registration ticks — a soothing, minimal 'technical drawing' feel."""
    from PIL import ImageDraw
    W, H = 2400, 1350
    rng = np.random.default_rng(7)
    yy = np.linspace(0, 1, H)[:, None]
    top = np.array([246, 242, 233.0]); bot = np.array([238, 232, 219.0])
    arr = np.repeat((top * (1 - yy) + bot * yy)[:, None, :], W, axis=1)
    arr += rng.normal(0, 1.3, (H, W, 1))           # faint grain (no banding)
    xg = np.linspace(-1, 1, W)[None, :]; yv = np.linspace(-1, 1, H)[:, None]
    arr *= np.clip(1 - (np.sqrt(xg ** 2 + yv ** 2) - 0.62) * 0.09, 0.94, 1.0)[..., None]
    img = Image.fromarray(np.clip(arr, 0, 255).astype(np.uint8), "RGB").convert("RGBA")
    ov = Image.new("RGBA", (W, H), (0, 0, 0, 0)); dr = ImageDraw.Draw(ov)
    warm = (150, 138, 116)
    for x in range(0, W, 80):                       # barely-there grid
        dr.line([(x, 0), (x, H)], fill=warm + (7,), width=1)
    for y in range(0, H, 80):
        dr.line([(0, y), (W, y)], fill=warm + (7,), width=1)
    m = 48
    dr.rectangle([m, m, W - m, H - m], outline=(120, 108, 88, 60), width=1)
    t = 22
    for cx, cy in [(m, m), (W - m, m), (m, H - m), (W - m, H - m)]:
        dr.line([(cx - t, cy), (cx + t, cy)], fill=(120, 108, 88, 130), width=1)
        dr.line([(cx, cy - t), (cx, cy + t)], fill=(120, 108, 88, 130), width=1)
    Image.alpha_composite(img, ov).convert("RGB").save(path)


def gen_drone_hero(path):
    """A blueprint / patent-style line-art quadcopter for the title hero."""
    _setup_sketch()
    ink = "#4A4338"; acc = "#5E7B72"; faint = "#A89E8B"
    fig = plt.figure(figsize=(8, 8)); ax = fig.add_axes([0, 0, 1, 1])
    ax.axis("off"); ax.set_aspect("equal"); ax.set_xlim(-11, 11); ax.set_ylim(-11.5, 11)
    motors = [(6.6, 6.6), (-6.6, 6.6), (6.6, -6.6), (-6.6, -6.6)]
    for mx, my in motors:                              # arms as twin-line tubes
        L = math.hypot(mx, my); px, py = -my / L, mx / L
        for s_ in (0.55, -0.55):
            ax.plot([px * s_, mx + px * s_], [py * s_, my + py * s_], color=ink, lw=1.6)
    for mx, my in motors:                              # prop discs (dashed)
        ax.add_patch(Circle((mx, my), 3.4, color=acc, alpha=0.07, lw=0))
        ax.add_patch(Circle((mx, my), 3.4, fill=False, ec=acc, lw=1.2, alpha=0.55,
                            linestyle=(0, (3, 3))))
        for ang in (40, 130):
            a = math.radians(ang)
            ax.plot([mx - 3.1 * math.cos(a), mx + 3.1 * math.cos(a)],
                    [my - 3.1 * math.sin(a), my + 3.1 * math.sin(a)], color=ink, lw=1.1, alpha=0.7)
    for mx, my in motors:                              # motor outlines
        ax.add_patch(Circle((mx, my), 1.5, fill=False, ec=ink, lw=1.6))
        ax.add_patch(Circle((mx, my), 0.45, fill=False, ec=ink, lw=1.1))
    ax.add_patch(FancyBboxPatch((-3.3, -2.7), 6.6, 5.4, boxstyle="round,pad=0,rounding_size=0.8",
                 fill=False, ec=ink, lw=2.0))
    ax.add_patch(FancyBboxPatch((-2.2, -1.7), 4.4, 3.4, boxstyle="round,pad=0,rounding_size=0.5",
                 fill=False, ec=ink, lw=1.0))
    ax.add_patch(FancyBboxPatch((-1.2, -4.4), 2.4, 1.9, boxstyle="round,pad=0,rounding_size=0.4",
                 fill=False, ec=ink, lw=1.5))           # gimbal
    ax.add_patch(Circle((0, -3.45), 0.5, fill=False, ec=ink, lw=1.2))
    ax.add_patch(Circle((0, -3.45), 0.16, fc=acc, ec="none"))
    ax.plot([-0.8, 0.8], [0, 0], color=faint, lw=0.9); ax.plot([0, 0], [-0.8, 0.8], color=faint, lw=0.9)
    ax.plot([6.6, 9.4], [6.6, 9.4], color=faint, lw=0.8)   # dimension leader
    ax.text(9.7, 9.7, "Ø PROP", color=faint, fontsize=10.5, ha="left", va="center", family="Consolas")
    fig.savefig(path, dpi=210, transparent=True, bbox_inches="tight", pad_inches=0.1)
    plt.close(fig)


# ===================================================================
#  HAND-DRAWN SKETCHES (matplotlib)
# ===================================================================
def _setup_sketch():
    matplotlib.rcParams.update({
        "path.sketch": None,                 # clean, professional lines (no wobble)
        "font.family": ["Segoe UI", "DejaVu Sans"],
        "savefig.transparent": True,
        "lines.solid_capstyle": "round",
    })


def _newax(w, h):
    fig = plt.figure(figsize=(w, h))
    ax = fig.add_axes([0, 0, 1, 1]); ax.axis("off"); ax.set_aspect("equal")
    return fig, ax


def _box(ax, x, y, w, h, label, ec=mINK, fc="none", fs=13, lw=2.6, tc=None, weight="bold"):
    ax.add_patch(FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0,rounding_size=1.1",
                                ec=ec, fc=fc, lw=lw, joinstyle="round"))
    if label:
        ax.text(x + w / 2, y + h / 2, label, ha="center", va="center",
                fontsize=fs, color=tc or ec, weight=weight, linespacing=1.2)


def _arr(ax, p1, p2, c=mINK, lw=2.6, style="-|>", rad=0.0, ls="-"):
    ax.add_patch(FancyArrowPatch(p1, p2, arrowstyle=style, mutation_scale=18, lw=lw,
                                 color=c, connectionstyle=f"arc3,rad={rad}",
                                 shrinkA=7, shrinkB=7, linestyle=ls))


def _txt(ax, x, y, t, c=mINK, fs=12, ha="center", va="center", rot=0, weight="bold"):
    ax.text(x, y, t, ha=ha, va=va, fontsize=fs, color=c, rotation=rot, weight=weight)


def _drone(ax, cx, cy, s, c=mINK):
    for dx, dy in [(-1, -1), (-1, 1), (1, -1), (1, 1)]:
        ex, ey = cx + dx * s, cy + dy * s
        ax.plot([cx, ex], [cy, ey], color=c, lw=2.6, solid_capstyle="round")
        ax.add_patch(Circle((ex, ey), s * 0.5, ec=c, fc="none", lw=2.4))
        ax.plot([ex - s * 0.5, ex + s * 0.5], [ey, ey], color=c, lw=1.4)
        ax.plot([ex, ex], [ey - s * 0.5, ey + s * 0.5], color=c, lw=1.4)
    ax.add_patch(Circle((cx, cy), s * 0.42, ec=c, fc="none", lw=2.8))


def _sat(ax, cx, cy, s, c=mINDIGO):
    ax.add_patch(Rectangle((cx - s * 0.45, cy - s * 0.45), s * 0.9, s * 0.9, ec=c, fc="none", lw=2.4))
    for sgn in (-1, 1):
        ax.plot([cx + sgn * s * 0.45, cx + sgn * s * 1.0], [cy, cy], color=c, lw=2.0)
        ax.add_patch(Rectangle((cx + sgn * s * 1.0 - (s * 0.18 if sgn > 0 else -s * 0.0),
                                cy - s * 0.55) if sgn > 0 else (cx - s * 1.4, cy - s * 0.55),
                               s * 0.4, s * 1.1, ec=c, fc="none", lw=2.0))
    for k in range(1, 4):
        ax.add_patch(matplotlib.patches.Arc((cx, cy - s * 0.6), s * k * 0.7, s * k * 0.5,
                     theta1=200, theta2=340, ec=c, lw=1.4))


def gen_sketches():
    _setup_sketch()

    # --- architecture ---
    fig, ax = _newax(9.0, 4.8); ax.set_xlim(0, 90); ax.set_ylim(0, 48)
    _box(ax, 4, 17, 17, 12, "Dashboard\n(React · Leaflet)", ec=mINDIGO, fs=12)
    _box(ax, 31, 6, 27, 36, "", ec=mINK, lw=2.8)
    _txt(ax, 44.5, 39, "Companion Computer", c=mINK, fs=13)
    _box(ax, 34, 30, 21, 6.5, "Trigger API", ec=mTERRA, fs=12)
    _box(ax, 34, 21.5, 21, 6.5, "Mission Executor", ec=mINK, fs=12)
    _box(ax, 34, 13, 21, 6.5, "Failsafe Monitor", ec=mOLIVE, fs=12)
    _box(ax, 67, 17, 18, 12, "ArduPilot", ec=mOLIVE, fs=12)
    _drone(ax, 76, 36, 2.6, c=mINK)
    _arr(ax, (21, 24), (31, 24), c=mINK, style="<|-|>")
    _txt(ax, 26, 27, "HTTP / WS", c=mINDIGO, fs=11)
    _arr(ax, (58, 24), (67, 24), c=mINK, style="<|-|>")
    _txt(ax, 62.5, 27, "MAVLink", c=mOLIVE, fs=11)
    fig.savefig(os.path.join(SK, "arch.png"), dpi=200, bbox_inches="tight", pad_inches=0.1)
    plt.close(fig)

    # --- navigation ---
    fig, ax = _newax(9.0, 5.0); ax.set_xlim(0, 90); ax.set_ylim(0, 50)
    for sx in (20, 45, 70):
        _sat(ax, sx, 43, 3.0)
        _arr(ax, (sx, 39), (45, 31), c=mINDIGO, lw=1.8, ls=(0, (4, 3)))
    _drone(ax, 45, 26, 3.4, c=mINK)
    _txt(ax, 45, 18.5, "DRONE — EKF: GPS · IMU · baro · compass", c=mINK, fs=12)
    _txt(ax, 45, 47.5, "GPS / GNSS satellites", c=mINDIGO, fs=12)
    _box(ax, 4, 22, 15, 9, "Companion\n(FastAPI)", ec=mTERRA, fs=11)
    _arr(ax, (19, 26.5), (37, 26.5), c=mINK, style="<|-|>")
    _txt(ax, 28, 29, "MAVLink", c=mTERRA, fs=10)
    _txt(ax, 81, 27, "✕", c=mTERRA, fs=22)
    _txt(ax, 81, 23, "target", c=mTERRA, fs=11)
    _arr(ax, (52, 25), (78, 26), c=mMUTE, lw=2.0, rad=-0.25, ls=(0, (2, 2)))
    _txt(ax, 65, 33, "lat/lon → goto", c=mMUTE, fs=10)
    fig.savefig(os.path.join(SK, "nav.png"), dpi=200, bbox_inches="tight", pad_inches=0.1)
    plt.close(fig)

    # --- state machine (snake) ---
    fig, ax = _newax(9.4, 4.6); ax.set_xlim(0, 94); ax.set_ylim(0, 46)
    r1 = ["IDLE", "CONNECTING", "WAITING_GPS", "ARMING", "TAKEOFF"]
    r2 = ["ENROUTE", "HOVERING", "RTL", "LANDED", "COMPLETED"]
    xs = [3, 21.5, 40, 58.5, 77]; w = 14.5; hb = 7
    for i, lab in enumerate(r1):
        _box(ax, xs[i], 31, w, hb, lab, ec=mINK, fs=10.5)
        if i:
            _arr(ax, (xs[i - 1] + w, 34.5), (xs[i], 34.5), c=mINK)
    for i, lab in enumerate(r2):
        c = mTERRA if lab in ("ENROUTE",) else (mOLIVE if lab == "COMPLETED" else mINK)
        _box(ax, xs[i], 14, w, hb, lab, ec=c, fs=10.5)
    _arr(ax, (xs[4] + w / 2, 31), (xs[4] + w / 2, 21), c=mINK, rad=0.0)
    for i in range(len(r2) - 1, 0, -1):
        _arr(ax, (xs[i], 17.5), (xs[i - 1] + w, 17.5), c=mINK)
    _box(ax, 24, 2, 16, 6.5, "ABORTED", ec=mTERRA, fs=10.5)
    _box(ax, 46, 2, 16, 6.5, "FAILED", ec=mTERRA, fs=10.5)
    _arr(ax, (40, 14), (36, 8.5), c=mTERRA, lw=1.8, ls=(0, (3, 3)))
    _txt(ax, 64, 6, "any phase →", c=mMUTE, fs=10, ha="left")
    fig.savefig(os.path.join(SK, "states.png"), dpi=200, bbox_inches="tight", pad_inches=0.1)
    plt.close(fig)

    # --- sequence (swimlane) ---
    fig, ax = _newax(8.6, 5.0); ax.set_xlim(0, 86); ax.set_ylim(0, 50)
    lanes = [("Client", 14), ("Companion / API", 45), ("Drone", 76)]
    for name, lx in lanes:
        _box(ax, lx - 11, 43, 22, 5, name, ec=mINK, fs=11)
        ax.plot([lx, lx], [3, 43], color=mMUTE, lw=1.6, ls=(0, (2, 3)))
    msgs = [
        (14, 45, 40, "POST /trigger", mTERRA, "-|>"),
        (45, 45, 35, "validate + queue", mMUTE, "self"),
        (45, 76, 29, "connect · arm · takeoff", mINK, "-|>"),
        (45, 76, 24, "goto target", mINK, "-|>"),
        (76, 45, 19, "mode confirmed", mOLIVE, "<|-"),
        (45, 76, 14, "RTL", mINK, "-|>"),
        (76, 45, 9, "landed", mOLIVE, "<|-"),
        (45, 14, 4.5, "telemetry / done", mINDIGO, "<|-"),
    ]
    for x1, x2, y, lab, c, st in msgs:
        if st == "self":
            ax.add_patch(matplotlib.patches.Arc((x1 + 6, y), 9, 4, theta1=-90, theta2=90, ec=c, lw=2.2))
            _arr(ax, (x1 + 6, y - 2), (x1 + 0.5, y - 2), c=c, lw=2.0)
            _txt(ax, x1 + 14, y, lab, c=c, fs=10, ha="left")
        else:
            _arr(ax, (x1, y), (x2, y), c=c, lw=2.3, style=st)
            _txt(ax, (x1 + x2) / 2, y + 1.8, lab, c=c, fs=10)
    fig.savefig(os.path.join(SK, "seq.png"), dpi=200, bbox_inches="tight", pad_inches=0.1)
    plt.close(fig)

    # --- failsafe tree ---
    fig, ax = _newax(9.0, 4.8); ax.set_xlim(0, 90); ax.set_ylim(0, 48)
    _drone(ax, 9, 24, 3.0, c=mINK); _txt(ax, 9, 16, "1 Hz monitor", c=mMUTE, fs=10)
    haz = [("Battery", 40), ("GPS fix", 31), ("Geofence", 22), ("Timeout", 13), ("Link", 4.5)]
    for lab, y in haz:
        _box(ax, 28, y, 15, 6.5, lab, ec=mINK, fs=11)
        _arr(ax, (17, 24), (28, y + 3.2), c=mMUTE, lw=1.8)

    def chip(x, y, lab, c):
        _box(ax, x, y, 11, 5, lab, ec=c, fs=10.5, tc=c)
    _arr(ax, (43, 43.2), (60, 45), c=mMUST, lw=1.8); chip(60, 42.5, "≤20% RTL", mMUST)
    _arr(ax, (43, 42), (60, 38), c=mTERRA, lw=1.8); chip(60, 35.5, "≤10% LAND", mTERRA)
    _arr(ax, (43, 34), (60, 31), c=mTERRA, lw=1.8); chip(60, 28.5, "lost → LAND", mTERRA)
    _arr(ax, (43, 25), (60, 22), c=mMUST, lw=1.8); chip(60, 19.5, "out → RTL", mMUST)
    _arr(ax, (43, 16), (60, 13), c=mMUST, lw=1.8); chip(60, 10.5, "over → RTL", mMUST)
    _arr(ax, (43, 7.5), (60, 4), c=mMUST, lw=1.8); chip(60, 1.5, "stale → RTL", mMUST)
    fig.savefig(os.path.join(SK, "failsafe.png"), dpi=200, bbox_inches="tight", pad_inches=0.1)
    plt.close(fig)

    # --- safety interlock (gates) ---
    fig, ax = _newax(9.4, 3.4); ax.set_xlim(0, 94); ax.set_ylim(0, 34)
    gates = [("Edge\nvalidation", mTERRA), ("Serial\nqueue", mMUST), ("GPS-lock\npre-arm", mINDIGO),
             ("Failsafe\narbiter", mOLIVE), ("Landing\ninterlock", mINK)]
    gx = 3; gw = 15; gap = 4
    for i, (lab, c) in enumerate(gates):
        _box(ax, gx, 9, gw, 14, lab, ec=c, fs=11)
        ax.add_patch(Circle((gx + 2.2, 21), 1.9, ec=c, fc="none", lw=2.4))
        _txt(ax, gx + 2.2, 21, str(i + 1), c=c, fs=11)
        if i:
            _arr(ax, (gx - gap, 16), (gx, 16), c=mINK, lw=2.6)
        gx += gw + gap
    fig.savefig(os.path.join(SK, "interlock.png"), dpi=200, bbox_inches="tight", pad_inches=0.1)
    plt.close(fig)

    # --- trajectory (hand-drawn map) ---
    fig, ax = _newax(7.4, 5.2); ax.set_xlim(0, 74); ax.set_ylim(0, 52)
    ax.add_patch(Rectangle((3, 3), 68, 46, ec=mMUTE, fc="none", lw=2.0))
    # north arrow
    _arr(ax, (66, 40), (66, 46), c=mINK, lw=2.0); _txt(ax, 66, 47.5, "N", c=mINK, fs=12)
    # home + target
    ax.add_patch(Circle((13, 12), 1.8, ec=mOLIVE, fc=mOLIVE, lw=2)); _txt(ax, 13, 8, "HOME", c=mOLIVE, fs=11)
    _txt(ax, 60, 42, "✕", c=mTERRA, fs=20); _txt(ax, 60, 38, "TARGET", c=mTERRA, fs=11)
    _arr(ax, (15, 14), (58, 40), c=mINK, lw=2.6, rad=-0.18)
    _txt(ax, 30, 33, "896 m @ 15 m, ~8 m/s", c=mINK, fs=10, rot=27)
    ax.add_patch(matplotlib.patches.Arc((59, 41), 7, 6, theta1=0, theta2=300, ec=mMUST, lw=2.2))
    _txt(ax, 59, 47, "hover", c=mMUST, fs=10)
    _arr(ax, (57, 39), (15, 14.5), c=mTERRA, lw=2.2, rad=-0.18, ls=(0, (4, 3)))
    _txt(ax, 40, 20, "RTL (return)", c=mTERRA, fs=10, rot=27)
    _txt(ax, 13, 4.6, "closest approach 0.4 m", c=mMUTE, fs=9, ha="left")
    fig.savefig(os.path.join(SK, "traj.png"), dpi=200, bbox_inches="tight", pad_inches=0.1)
    plt.close(fig)

    # --- full system (hardware + software) ---
    fig, ax = _newax(9.6, 5.2); ax.set_xlim(0, 96); ax.set_ylim(0, 52)
    ax.add_patch(FancyBboxPatch((4, 22), 88, 27, boxstyle="round,pad=0,rounding_size=1.5",
                 ec=mMUTE, fc="none", lw=1.8, linestyle=(0, (5, 4))))
    ax.add_patch(FancyBboxPatch((4, 2), 88, 16, boxstyle="round,pad=0,rounding_size=1.5",
                 ec=mMUTE, fc="none", lw=1.8, linestyle=(0, (5, 4))))
    _txt(ax, 13, 47, "ON THE DRONE", c=mMUTE, fs=11, ha="left")
    _txt(ax, 14, 15.5, "GROUND STATION", c=mMUTE, fs=11, ha="left")
    _box(ax, 8, 35, 18, 9, "Sensors\nGNSS·IMU·baro", ec=mINDIGO, fs=9.5)
    _box(ax, 33, 35, 18, 9, "Flight\nController", ec=mINK, fs=10)
    _box(ax, 58, 31, 31, 15, "", ec=mTERRA, lw=2.6)
    _txt(ax, 73.5, 43, "Companion Computer (Pi)", c=mTERRA, fs=10)
    for j, lab in enumerate(["Trigger API", "Mission Executor", "Failsafe Monitor"]):
        _txt(ax, 73.5, 39 - j * 3, lab, c=mINK, fs=9)
    _box(ax, 33, 24, 18, 7, "ESC + Motors", ec=mOLIVE, fs=9.5)
    _box(ax, 8, 24, 16, 7, "SiK radio", ec=mMUST, fs=9.5)
    _box(ax, 8, 5, 16, 8, "SiK radio", ec=mMUST, fs=9.5)
    _box(ax, 33, 4, 34, 10, "Ground laptop\n(GCS · Dashboard)", ec=mINDIGO, fs=10)
    _arr(ax, (26, 39.5), (33, 39.5), c=mINK); _txt(ax, 29.5, 41.6, "data", c=mINDIGO, fs=8.5)
    _arr(ax, (51, 39.5), (58, 39.5), c=mINK, style="<|-|>"); _txt(ax, 54.5, 41.6, "UART", c=mTERRA, fs=8.5)
    _arr(ax, (42, 35), (42, 31), c=mINK); _txt(ax, 46.5, 33, "PWM", c=mOLIVE, fs=8.5)
    _arr(ax, (33, 37), (24, 30.5), c=mINK)
    _arr(ax, (16, 24), (16, 13), c=mMUST, style="<|-|>", ls=(0, (4, 3)))
    _txt(ax, 24.5, 18.5, "433 MHz RF", c=mMUST, fs=9, ha="left")
    _arr(ax, (24, 9), (33, 9), c=mINK, style="<|-|>"); _txt(ax, 28.5, 11, "USB", c=mINDIGO, fs=8.5)
    fig.savefig(os.path.join(SK, "system.png"), dpi=200, bbox_inches="tight", pad_inches=0.1)
    plt.close(fig)

    # --- methodology ---
    fig, ax = _newax(9.4, 3.4); ax.set_xlim(0, 94); ax.set_ylim(0, 34)
    stg = [("Requirements", mINDIGO), ("Design", mTERRA), ("Implement", mINK),
           ("SITL\nsimulate", mOLIVE), ("Verify &\nharden", mMUST)]
    sx = 3; sw = 15; cxs = []
    for i, (lab, c) in enumerate(stg):
        _box(ax, sx, 16, sw, 11, lab, ec=c, fs=10.5)
        cxs.append((sx, sx + sw))
        if i:
            _arr(ax, (cxs[i - 1][1], 21.5), (sx, 21.5), c=mINK, lw=2.4)
        sx += sw + 4
    _arr(ax, (cxs[4][0] + sw / 2, 16), (cxs[2][0] + sw / 2, 16), c=mTERRA, lw=2.0, rad=0.55, ls=(0, (4, 3)))
    _txt(ax, (cxs[2][1] + cxs[4][0]) / 2 + 2, 5.5, "iterate", c=mTERRA, fs=10.5)
    fig.savefig(os.path.join(SK, "methodology.png"), dpi=200, bbox_inches="tight", pad_inches=0.1)
    plt.close(fig)

    # --- working / dataflow ---
    fig, ax = _newax(9.6, 4.8); ax.set_xlim(0, 96); ax.set_ylim(0, 48)
    _box(ax, 4, 30, 16, 9, "Operator\n(HTTP)", ec=mINDIGO, fs=10)
    _box(ax, 24, 30, 21, 9, "Trigger API\nvalidate·geofence·queue", ec=mTERRA, fs=8.5)
    _box(ax, 49, 26, 24, 14, "", ec=mINK, lw=2.6)
    _txt(ax, 61, 36, "Mission Executor", c=mINK, fs=10.5)
    _txt(ax, 61, 30.3, "connect·arm·takeoff·\ngoto·hover·RTL·land", c=mMUTE, fs=8)
    _box(ax, 77, 30, 16, 9, "Autopilot\nGUIDED", ec=mOLIVE, fs=10)
    _box(ax, 77, 16, 16, 8, "GNSS + EKF", ec=mINDIGO, fs=9.5)
    _box(ax, 49, 13, 24, 8, "Failsafe monitor (1 Hz)", ec=mOLIVE, fs=9)
    _box(ax, 21, 13, 22, 9, "Dashboard\n(live telemetry)", ec=mINDIGO, fs=9.5)
    _arr(ax, (20, 34.5), (24, 34.5), c=mINK)
    _arr(ax, (45, 34.5), (49, 34.5), c=mINK)
    _arr(ax, (73, 34), (77, 34), c=mINK, style="<|-|>")
    _arr(ax, (85, 30), (85, 24), c=mINK, style="<|-|>"); _txt(ax, 89.5, 27, "pos", c=mINDIGO, fs=8)
    _arr(ax, (61, 26), (61, 21), c=mOLIVE, style="<|-|>"); _txt(ax, 65, 23.5, "watch", c=mOLIVE, fs=8)
    _arr(ax, (77, 31), (43, 17.5), c=mMUTE, lw=1.8, ls=(0, (3, 3)), rad=0.12)
    _txt(ax, 60, 14.5, "telemetry", c=mMUTE, fs=9)
    fig.savefig(os.path.join(SK, "working.png"), dpi=200, bbox_inches="tight", pad_inches=0.1)
    plt.close(fig)

    # --- how coordinates are received ---
    fig, ax = _newax(7.8, 5.4); ax.set_xlim(0, 78); ax.set_ylim(0, 54)
    steps = [("POST /trigger  {lat, lon, alt, hover}", mTERRA),
             ("Validate bounds  (-90..90, -180..180)", mINK),
             ("Geofence  <= 5 km from home ?", mINK),
             ("Enqueue MissionSpec  (priority)", mINK),
             ("Connect - GPS lock - arm - takeoff", mINK),
             ("simple_goto(lat, lon, alt)  ->  GUIDED", mINDIGO),
             ("Autopilot navigates by EKF (GPS)", mOLIVE),
             ("distance <= 5 m ?   ->  ARRIVED", mOLIVE)]
    y = 47; bx = 6; bw = 46; bh = 4.6; ys = []
    for i, (lab, c) in enumerate(steps):
        _box(ax, bx, y, bw, bh, lab, ec=c, fs=9)
        ys.append(y)
        if i:
            _arr(ax, (bx + bw / 2, ys[i - 1]), (bx + bw / 2, y + bh), c=mINK, lw=2.2)
        y -= 6.3
    _box(ax, 58, ys[2], 16, bh, "HTTP 400\nreject", ec=mTERRA, fs=8.5)
    _arr(ax, (bx + bw, ys[2] + bh / 2), (58, ys[2] + bh / 2), c=mTERRA, lw=2.0)
    _arr(ax, (bx + bw, ys[7] + bh / 2), (bx + bw + 5, ys[6] + bh / 2), c=mOLIVE, lw=1.8, rad=-0.5)
    _txt(ax, bx + bw + 9, (ys[6] + ys[7]) / 2 + 2, "loop", c=mOLIVE, fs=8.5, ha="left")
    fig.savefig(os.path.join(SK, "coords.png"), dpi=200, bbox_inches="tight", pad_inches=0.1)
    plt.close(fig)

    # --- how it reaches the destination (control loop) ---
    fig, ax = _newax(8.6, 4.6); ax.set_xlim(0, 86); ax.set_ylim(0, 46)
    _box(ax, 30, 36, 26, 7, "Target  (lat, lon)", ec=mTERRA, fs=10.5)
    _box(ax, 56, 19, 26, 9, "Compare -> error\n(current to target)", ec=mINK, fs=9.5)
    _box(ax, 30, 3, 26, 7, "Autopilot GUIDED steers", ec=mOLIVE, fs=10)
    _box(ax, 4, 19, 22, 9, "Drone moves\nnew EKF position", ec=mINDIGO, fs=9.5)
    _arr(ax, (50, 36), (66, 28), c=mINK, rad=-0.2)
    _arr(ax, (66, 19), (50, 10.5), c=mINK, rad=-0.2)
    _arr(ax, (30, 6.5), (20, 19), c=mINK, rad=-0.2)
    _arr(ax, (15, 28), (34, 36), c=mINK, rad=-0.2)
    _txt(ax, 14, 33, "feedback", c=mINDIGO, fs=9, ha="left")
    _txt(ax, 71, 32, "<= 5 m -> ARRIVED", c=mOLIVE, fs=9.5)
    _arr(ax, (70, 27), (72, 30.5), c=mOLIVE, lw=1.8)
    _txt(ax, 71, 11, "no progress -> ABORT", c=mTERRA, fs=9.5)
    _arr(ax, (70, 20), (72, 13), c=mTERRA, lw=1.8)
    fig.savefig(os.path.join(SK, "reach.png"), dpi=200, bbox_inches="tight", pad_inches=0.1)
    plt.close(fig)

    # --- obstacle avoidance: map-based detour (implemented) + roadmap ---
    fig, ax = _newax(9.2, 4.8); ax.set_xlim(0, 92); ax.set_ylim(0, 48)
    ax.add_patch(FancyBboxPatch((4, 3), 60, 42, boxstyle="round,pad=0,rounding_size=2",
                 ec=mMUTE, fc="none", lw=1.6, linestyle=(0, (5, 4))))
    _txt(ax, 8, 42.5, "GEOFENCE", c=mMUTE, fs=9.5, ha="left")
    ax.add_patch(Circle((33, 24), 9, ec=mTERRA, fc="none", lw=2.2, linestyle=(0, (4, 3))))
    ax.add_patch(Circle((33, 24), 12.5, ec=mMUST, fc="none", lw=1.3, linestyle=(0, (2, 3))))
    _txt(ax, 33, 24, "keep-out\nzone", c=mTERRA, fs=9)
    _txt(ax, 33, 38.5, "+ clearance", c=mMUST, fs=8.5)
    ax.add_patch(Circle((9, 12), 1.6, ec=mOLIVE, fc=mOLIVE)); _txt(ax, 9, 8, "HOME", c=mOLIVE, fs=9)
    _txt(ax, 59.5, 31.5, "X", c=mTERRA, fs=16); _txt(ax, 59, 27.5, "TARGET", c=mTERRA, fs=9)
    ax.plot([9, 58], [12, 31], color=mMUTE, lw=1.5, ls=(0, (2, 2)))
    _txt(ax, 28, 13.5, "direct = blocked", c=mMUTE, fs=8.5, rot=20)
    _arr(ax, (9, 13.5), (24, 39), c=mINK, lw=2.4, rad=-0.15)
    _arr(ax, (24, 39), (45, 37), c=mINK, lw=2.4)
    _arr(ax, (45, 37), (58, 31.5), c=mINK, lw=2.4, rad=-0.12)
    ax.add_patch(Circle((24, 39), 1.2, ec=mINK, fc=mINK))
    ax.add_patch(Circle((45, 37), 1.2, ec=mINK, fc=mINK))
    _txt(ax, 33, 45.8, "map-based detour (implemented)", c=mINK, fs=10)
    ax.add_patch(FancyBboxPatch((68, 6), 22, 36, boxstyle="round,pad=0,rounding_size=2",
                 ec=mMUTE, fc="none", lw=1.6, linestyle=(0, (3, 3))))
    _txt(ax, 79, 39, "FUTURE (roadmap)", c=mMUTE, fs=9.5)
    _drone(ax, 74, 28, 2.0, c=mMUTE)
    ax.plot([76, 86, 86, 76], [28, 31.5, 24.5, 28], color=mTERRA, lw=1.6)
    ax.add_patch(Rectangle((86, 24.5), 3, 7, ec=mINK, fc="none", lw=1.8))
    _txt(ax, 79, 13, "sensor-based\nreactive avoidance", c=mMUTE, fs=8.5)
    fig.savefig(os.path.join(SK, "obstacles.png"), dpi=200, bbox_inches="tight", pad_inches=0.1)
    plt.close(fig)

    # --- hardware wiring ---
    fig, ax = _newax(9.4, 5.0); ax.set_xlim(0, 94); ax.set_ylim(0, 50)
    _box(ax, 37, 19, 22, 15, "", ec=mINK, lw=2.8)
    _txt(ax, 48, 30, "Flight Controller", c=mINK, fs=11)
    _txt(ax, 48, 26.5, "(Pixhawk / Cube)", c=mMUTE, fs=8.5)
    _txt(ax, 58.5, 31, "TELEM1", c=mTERRA, fs=7.5, ha="right")
    _txt(ax, 58.5, 22.5, "UART", c=mINDIGO, fs=7.5, ha="right")
    _txt(ax, 37.5, 31, "GPS", c=mOLIVE, fs=7.5, ha="left")
    _txt(ax, 37.5, 22.5, "PWR", c=mMUST, fs=7.5, ha="left")
    _box(ax, 68, 36, 22, 9, "SiK radio  433 MHz", ec=mTERRA, fs=9)
    _box(ax, 68, 12, 22, 9, "Companion Pi", ec=mINDIGO, fs=9)
    _box(ax, 4, 36, 22, 9, "GPS + compass", ec=mOLIVE, fs=9)
    _box(ax, 4, 12, 22, 9, "Power module", ec=mMUST, fs=9)
    _box(ax, 37, 3, 22, 8, "ESC + Motors x4", ec=mINK, fs=9)
    _arr(ax, (59, 31), (68, 40.5), c=mTERRA, lw=2.0, style="<|-|>")
    _txt(ax, 65.5, 36.5, "6-wire JST-GH", c=mTERRA, fs=7, ha="left")
    _arr(ax, (59, 22.5), (68, 16.5), c=mINDIGO, lw=2.0, style="<|-|>")
    _txt(ax, 64.5, 21, "TX·RX·GND", c=mINDIGO, fs=7, ha="left")
    _arr(ax, (37, 31), (26, 40.5), c=mOLIVE, lw=2.0)
    _arr(ax, (37, 22.5), (26, 16.5), c=mMUST, lw=2.0)
    _arr(ax, (48, 19), (48, 11), c=mINK, lw=2.0); _txt(ax, 51, 15, "PWM", c=mINK, fs=7, ha="left")
    _txt(ax, 48, 47, "FC TX -> radio RX   (cross TX / RX)", c=mMUTE, fs=9)
    fig.savefig(os.path.join(SK, "wiring.png"), dpi=200, bbox_inches="tight", pad_inches=0.1)
    plt.close(fig)


gen_paper(PAPER_PNG)
gen_drone_hero(DRONE_PNG)
gen_sketches()

# ===================================================================
#  DECK
# ===================================================================
prs = Presentation(); prs.slide_width = SW; prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def soft_shadow(shape, blur=0.10, dist=0.05, alpha=26, color="000000"):
    spPr = shape._element.spPr
    for e in spPr.findall(qn("a:effectLst")):
        spPr.remove(e)
    spPr.append(parse_xml(
        f'<a:effectLst {nsdecls("a")}><a:outerShdw blurRad="{int(blur*914400)}" '
        f'dist="{int(dist*914400)}" dir="5400000" rotWithShape="0">'
        f'<a:srgbClr val="{color}"><a:alpha val="{int(alpha*1000)}"/></a:srgbClr>'
        f'</a:outerShdw></a:effectLst>'))


def _alpha(srgb, pct):
    for e in srgb.findall(qn("a:alpha")):
        srgb.remove(e)
    srgb.append(parse_xml(f'<a:alpha {nsdecls("a")} val="{int(pct*1000)}"/>'))


def fill_alpha(shape, pct):
    sF = shape._element.spPr.find(qn("a:solidFill"))
    if sF is not None:
        c = sF.find(qn("a:srgbClr"))
        if c is not None:
            _alpha(c, pct)


def line_alpha(shape, pct):
    ln = shape._element.spPr.find(qn("a:ln"))
    if ln is not None:
        sF = ln.find(qn("a:solidFill"))
        if sF is not None:
            c = sF.find(qn("a:srgbClr"))
            if c is not None:
                _alpha(c, pct)


_idx = [0]


def slide():
    _idx[0] += 1
    s = prs.slides.add_slide(BLANK)
    s.shapes.add_picture(PAPER_PNG, 0, 0, SW, SH)
    return s


def rect(s, x, y, w, h, fill=None, line=None, lw=1.0, rounded=False, rad=0.12, shadow=False):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, x, y, w, h)
    if rounded:
        short = min(int(w), int(h))
        if short > 0:
            try:
                shp.adjustments[0] = max(0.0, min(0.5, float(Inches(rad)) / short))
            except Exception:
                pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    if shadow:
        soft_shadow(shp)
    return shp


def glass(s, x, y, w, h, fa=92, ba=100, border=HAIR, rad=0.05, shadow=False, hi=False):
    """A flat, crisp technical panel — warm off-white with a clean hairline
    (or muted-accent) border. No glossy effects; reads like a drawing frame."""
    shp = rect(s, x, y, w, h, fill=PANEL, line=border, lw=1.1, rounded=True, rad=rad)
    fill_alpha(shp, fa); line_alpha(shp, max(ba, 72))
    if shadow:
        soft_shadow(shp, blur=0.10, dist=0.04, alpha=12, color="6B5A3A")
    return shp


def text(s, x, y, w, h, runs, size=18, color=INK, font=F_BODY, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0, italic=False):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor; tf.auto_size = MSO_AUTO_SIZE.NONE
    if isinstance(runs, str):
        runs = [runs]
    for i, ln in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = spacing; p.space_after = Pt(2)
        txt, opt = ln if isinstance(ln, tuple) else (ln, {})
        run = p.add_run(); run.text = txt; f = run.font
        f.size = Pt(opt.get("size", size)); f.name = opt.get("font", font)
        f.bold = opt.get("bold", bold); f.italic = opt.get("italic", italic)
        f.color.rgb = opt.get("color", color)
    return tb


def header(s, kicker, title, color=INDIGO):
    rect(s, Inches(0.74), Inches(0.62), Inches(0.13), Inches(0.13), fill=color)
    text(s, Inches(1.0), Inches(0.55), Inches(10), Inches(0.35),
         kicker.upper(), size=12, color=color, font=F_MONO, bold=True)
    text(s, Inches(0.7), Inches(0.92), Inches(12), Inches(0.85),
         title, size=31, color=INK, font=F_TITLE, bold=True)
    rect(s, Inches(0.74), Inches(1.66), Inches(11.85), Pt(1.1), fill=HAIR)
    rect(s, Inches(0.74), Inches(1.60), Pt(1.6), Inches(0.13), fill=color)        # end ticks
    rect(s, Inches(12.55), Inches(1.60), Pt(1.6), Inches(0.13), fill=color)


def footer(s, n=None):
    rect(s, Inches(0.74), Inches(6.98), Inches(11.85), Pt(0.9), fill=HAIR)
    text(s, Inches(0.74), Inches(7.06), Inches(8), Inches(0.3),
         "DRONE SAFETY SYSTEM", size=9.5, color=MUTED, font=F_MONO)
    text(s, Inches(10.95), Inches(7.06), Inches(1.65), Inches(0.3),
         f"{_idx[0]:02d} / 24", size=9.5, color=MUTED, font=F_MONO, align=PP_ALIGN.RIGHT)


def bullets(s, x, y, w, h, items, size=16, gap=9, color=INK, mcolor=TERRA):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.line_spacing = 1.1
        r1 = p.add_run(); r1.text = "•  "
        r1.font.size = Pt(size); r1.font.name = F_TITLE; r1.font.bold = True
        r1.font.color.rgb = mcolor
        if isinstance(it, tuple):
            lead, rest = it
            r2 = p.add_run(); r2.text = lead
            r2.font.size = Pt(size); r2.font.name = F_TITLE; r2.font.bold = True
            r2.font.color.rgb = color
            r3 = p.add_run(); r3.text = rest
            r3.font.size = Pt(size); r3.font.name = F_BODY; r3.font.color.rgb = SOFT
        else:
            r2 = p.add_run(); r2.text = it
            r2.font.size = Pt(size); r2.font.name = F_BODY; r2.font.color.rgb = color
    return tb


def sketch(s, name, x, y, w, h, caption=None):
    path = os.path.join(SK, name)
    iw, ih = Image.open(path).size
    bw = int(w); bh = int(h) - (int(Inches(0.4)) if caption else 0)
    scale = min(bw / iw, bh / ih)
    nw, nh = int(iw * scale), int(ih * scale)
    px = int(x) + (int(w) - nw) // 2; py = int(y) + (bh - nh) // 2
    s.shapes.add_picture(path, Emu(px), Emu(py), Emu(nw), Emu(nh))
    if caption:
        text(s, x, Emu(int(y) + bh), w, Inches(0.4), caption, size=11, color=MUTED,
             font=F_MONO, italic=False, align=PP_ALIGN.CENTER)


# ---------------- 1 · TITLE ----------------
s = slide()
# blueprint line-art drone, cut off on the top-right edge
s.shapes.add_picture(DRONE_PNG, Emu(int(SW) - int(Inches(5.0))), Inches(-0.5),
                     Inches(7.7), Inches(7.7))
text(s, Inches(0.9), Inches(0.62), Inches(7), Inches(0.3),
     "REV 1.1  /  SITL-VERIFIED  /  2026", size=10.5, color=MUTED, font=F_MONO)
text(s, Inches(0.9), Inches(1.55), Inches(7.4), Inches(0.4),
     "AUTONOMOUS UAV  /  COMPANION-COMPUTER SAFETY",
     size=12, color=INDIGO, font=F_MONO, bold=True)
text(s, Inches(0.85), Inches(2.05), Inches(6.7), Inches(2.2),
     ["Drone Safety", "System"], size=58, color=INK, font=F_HERO, bold=True, spacing=0.94)
rect(s, Inches(0.92), Inches(4.62), Inches(2.2), Pt(2.4), fill=INDIGO)
text(s, Inches(0.9), Inches(4.84), Inches(6.4), Inches(0.9),
     "Verified dispatch & failsafe arbitration for fully autonomous flight.",
     size=18, color=SOFT, font=F_BODY, spacing=1.15)
kpis = [("0.4 m", "TERMINAL ACCURACY"), ("8 / 8", "ACCEPTANCE CHECKS"),
        ("54", "UNIT TESTS"), ("2", "PATENTS")]
x = Inches(0.9)
for i, (val, lab) in enumerate(kpis):
    text(s, x, Inches(5.98), Inches(1.55), Inches(0.6), val, size=28, color=INK,
         font=F_HERO, bold=True)
    text(s, x, Inches(6.64), Inches(1.6), Inches(0.3), lab, size=8.5, color=MUTED, font=F_MONO)
    if i < 3:
        rect(s, Emu(int(x) + int(Inches(1.52))), Inches(6.04), Pt(1.0), Inches(0.76), fill=HAIR)
    x = Emu(int(x) + int(Inches(1.64)))
text(s, Inches(0.9), Inches(7.08), Inches(8), Inches(0.3),
     "github.com/SV-1411/drone", size=11.5, color=MUTED, font=F_MONO)

# ---------------- 2 · WHAT IT IS ----------------
s = slide(); header(s, "Overview", "What this system is")
text(s, Inches(0.74), Inches(1.9), Inches(11.85), Inches(1.0),
     "A companion-computer stack that turns one network trigger into a complete autonomous "
     "mission — fly to a GPS target, hover, and return — with no pilot, and aviation-grade "
     "safety woven through every phase.", size=19, color=SOFT, spacing=1.18)
cards = [
    ("Trigger-to-flight", "POST a coordinate; the drone arms, takes off, navigates, hovers and returns on its own.", TERRA),
    ("Verified commands", "Every flight-mode change is confirmed from the autopilot's own telemetry, never assumed.", OLIVE),
    ("Arbitrated failsafes", "Battery, GPS, geofence & timeout hazards resolve to one debounced action.", INDIGO),
    ("Provable safety", "A landing interlock makes it impossible to start a mission on an airborne drone.", MUST),
]
x = Inches(0.74); w = Inches(2.83)
for t, d, c in cards:
    glass(s, x, Inches(3.05), w, Inches(2.95), border=c, ba=30)
    rect(s, Emu(int(x) + int(Inches(0.22))), Inches(3.3),
         Emu(int(w) - int(Inches(0.44))), Pt(2.2), fill=c)
    text(s, Emu(int(x) + int(Inches(0.26))), Inches(3.5), Emu(int(w) - int(Inches(0.52))),
         Inches(0.9), t, size=18, color=c, font=F_TITLE, bold=True)
    text(s, Emu(int(x) + int(Inches(0.26))), Inches(4.35), Emu(int(w) - int(Inches(0.52))),
         Inches(1.5), d, size=14.5, color=SOFT, spacing=1.14)
    x = Emu(int(x) + int(Inches(3.0)))
footer(s, 2)

# ---------------- 3 · PROBLEM ----------------
s = slide(); header(s, "Motivation", "Autonomy removes the operator who catches failures")
bullets(s, Inches(0.74), Inches(2.0), Inches(6.0), Inches(4.5), [
    ("The autopilot can lie. ", "A mode-change is reported accepted while the autopilot silently stays in its old mode — seen on real ArduCopter."),
    ("No human is watching. ", "In autonomous dispatch nobody notices a dropped command — and the same path carries emergency RTL/LAND."),
    ("Missions can overlap. ", "Treating mission N's end as license to start N+1 can arm a still-airborne aircraft — undefined & hazardous."),
    ("Sensors glitch. ", "One bad GPS sample, taken naively, triggers an instant land-in-place — over a road, water, or a crowd."),
], size=16.5, gap=15)
glass(s, Inches(7.15), Inches(2.05), Inches(5.45), Inches(4.25), border=TERRA, ba=28)
rect(s, Inches(7.15), Inches(2.05), Inches(0.12), Inches(4.25), fill=TERRA, rounded=False)
text(s, Inches(7.55), Inches(2.35), Inches(4.8), Inches(0.4),
     "The core insight", size=15, color=TERRA, font=F_TITLE, italic=True)
text(s, Inches(7.55), Inches(2.85), Inches(4.75), Inches(1.4),
     "“A command issued is not a command executed.”",
     size=24, color=INK, font=F_TITLE, bold=True, spacing=1.1)
text(s, Inches(7.55), Inches(4.45), Inches(4.75), Inches(1.8),
     "So we treat the autopilot as untrusted: demand telemetry-confirmed proof of every "
     "transition — then guarantee the queue can never hand off an airborne drone.",
     size=16, color=SOFT, spacing=1.18, italic=True)
footer(s, 3)

# ---------------- 4 · ARCHITECTURE ----------------
s = slide(); header(s, "System architecture", "Four cleanly separated parts")
sketch(s, "arch.png", Inches(0.6), Inches(1.85), Inches(8.4), Inches(5.0),
       caption="Figure 1 · System architecture")
glass(s, Inches(9.15), Inches(2.0), Inches(3.5), Inches(4.5), border=INDIGO, ba=28)
bullets(s, Inches(9.4), Inches(2.25), Inches(3.05), Inches(4.1), [
    ("Trigger API ", "— validate, queue, stream."),
    ("Executor ", "— the 12-state flight machine."),
    ("Failsafe monitor ", "— 1 Hz hazard arbiter."),
    ("Dashboard ", "— live Leaflet map."),
], size=14, gap=12, mcolor=INDIGO)
footer(s, 4)

# ---------------- SYSTEM DIAGRAM ----------------
s = slide(); header(s, "System", "The full system — hardware & software", color=INDIGO)
sketch(s, "system.png", Inches(0.55), Inches(1.85), Inches(8.6), Inches(5.0),
       caption="Figure 2 · End-to-end system (drone + ground)")
glass(s, Inches(9.3), Inches(2.0), Inches(3.35), Inches(4.5), border=INDIGO, ba=28)
bullets(s, Inches(9.52), Inches(2.25), Inches(2.95), Inches(4.1), [
    ("On the drone ", "— sensors → flight controller → companion Pi → motors."),
    ("Companion ", "— runs the Trigger API, executor & failsafe."),
    ("Link ", "— SiK 433 MHz radio to the ground laptop."),
    ("Ground ", "— GCS / dashboard for monitoring & override."),
], size=13, gap=11, mcolor=INDIGO)
footer(s)

# ---------------- HARDWARE WIRING ----------------
s = slide(); header(s, "Hardware", "How the radio & companion are wired in", color=TERRA)
sketch(s, "wiring.png", Inches(0.55), Inches(1.85), Inches(8.6), Inches(5.0),
       caption="Figure · Flight-controller wiring (SiK 433 + companion)")
glass(s, Inches(9.3), Inches(2.0), Inches(3.35), Inches(4.5), border=TERRA, ba=28)
bullets(s, Inches(9.52), Inches(2.25), Inches(2.95), Inches(4.1), [
    ("TELEM1 → SiK radio ", "— 6-wire JST-GH; cross FC TX → radio RX."),
    ("UART → Companion ", "— TX·RX·GND; runs the autonomy stack."),
    ("GPS & power ", "— to their dedicated ports."),
    ("MAIN OUT → ESCs ", "— PWM to the four motors."),
], size=13, gap=11)
footer(s)

# ---------------- METHODOLOGY ----------------
s = slide(); header(s, "Methodology", "How the system was built & proven")
text(s, Inches(0.7), Inches(4.78), Inches(11.9), Inches(0.35),
     "Figure 3 · Iterative, test-driven development method", size=13, color=MUTED,
     font=F_BODY, italic=True, align=PP_ALIGN.CENTER)
sketch(s, "methodology.png", Inches(0.7), Inches(1.95), Inches(11.9), Inches(2.75))
glass(s, Inches(0.74), Inches(5.15), Inches(11.85), Inches(1.4), border=TERRA, ba=26)
bullets(s, Inches(1.05), Inches(5.3), Inches(11.2), Inches(1.2), [
    ("Test-driven ", "— design → implement → simulate in SITL → verify with 37 unit + E2E tests → harden, repeat."),
    ("Evidence-led ", "— every safety claim is pinned by an automated test before it enters the patent drafts."),
], size=14, gap=8)
footer(s)

# ---------------- WORKING / DATAFLOW ----------------
s = slide(); header(s, "Working", "End-to-end: how it all flows", color=OLIVE)
sketch(s, "working.png", Inches(0.55), Inches(1.85), Inches(8.7), Inches(5.0),
       caption="Figure 4 · Working & data flow, input to flight")
glass(s, Inches(9.35), Inches(2.0), Inches(3.3), Inches(4.5), border=OLIVE, ba=28)
bullets(s, Inches(9.57), Inches(2.25), Inches(2.9), Inches(4.1), [
    ("Input ", "— operator POSTs a coordinate."),
    ("Process ", "— API validates & queues; executor flies the phases."),
    ("Control ", "— autopilot in GUIDED on GNSS-fused position."),
    ("Feedback ", "— live telemetry out; failsafe watches throughout."),
], size=13, gap=11, mcolor=OLIVE)
footer(s)

# ---------------- 5 · NAVIGATION ----------------
s = slide(); header(s, "Navigation", "How it knows where it is: GPS, not cell towers")
sketch(s, "nav.png", Inches(0.55), Inches(1.8), Inches(8.4), Inches(4.55),
       caption="Figure 2 · Positioning by satellite GNSS + sensor fusion")
glass(s, Inches(9.15), Inches(1.95), Inches(3.5), Inches(3.0), border=INDIGO, ba=28)
bullets(s, Inches(9.4), Inches(2.2), Inches(3.05), Inches(2.6), [
    ("Positioning ", "— GNSS + IMU + baro + compass, fused onboard."),
    ("Data link ", "— MAVLink over radio / WiFi / LTE."),
], size=14, gap=11, mcolor=INDIGO)
glass(s, Inches(0.74), Inches(6.5), Inches(11.85), Inches(0.62), border=TERRA, ba=26, hi=False)
text(s, Inches(1.0), Inches(6.55), Inches(11.3), Inches(0.5),
     [("Cell / 5G is only an optional link — never how the drone locates itself.  ",
       {"bold": True, "color": INK, "font": F_TITLE, "size": 14}),
      ("Positioning is satellite GNSS.", {"color": SOFT, "size": 14, "italic": True})],
     anchor=MSO_ANCHOR.MIDDLE)
footer(s, 5)

# ---------------- HOW IT RECEIVES COORDINATES ----------------
s = slide(); header(s, "Input", "How the drone receives coordinates", color=TERRA)
sketch(s, "coords.png", Inches(0.6), Inches(1.85), Inches(7.4), Inches(5.0),
       caption="Figure 6 · From HTTP request to a navigated target")
glass(s, Inches(8.15), Inches(2.0), Inches(4.45), Inches(4.5), border=TERRA, ba=28)
bullets(s, Inches(8.4), Inches(2.25), Inches(4.0), Inches(4.1), [
    ("Coordinate in ", "— one POST with lat/lon (+ optional alt, hover)."),
    ("Checked twice ", "— value bounds, then geofence vs home (else 400)."),
    ("Queued ", "— as a MissionSpec; missions run strictly serially."),
    ("Handed over ", "— as a GUIDED-mode goto to the autopilot."),
    ("Closed-loop ", "— flies until within 5 m (haversine distance)."),
], size=14, gap=11)
footer(s)

# ---------------- HOW IT REACHES THE DESTINATION ----------------
s = slide(); header(s, "Guidance", "How it reaches the destination", color=INDIGO)
sketch(s, "reach.png", Inches(0.6), Inches(1.95), Inches(7.6), Inches(4.7),
       caption="Figure 7 · The goto control loop")
glass(s, Inches(8.35), Inches(2.0), Inches(4.25), Inches(4.5), border=INDIGO, ba=28)
bullets(s, Inches(8.6), Inches(2.25), Inches(3.75), Inches(4.1), [
    ("Closed loop ", "— compares EKF position to target every second."),
    ("Autopilot steers ", "— GUIDED mode corrects continuously."),
    ("Arrival ", "— declared within 5 m tolerance (RTK → ~0.3 m)."),
    ("Stall guard ", "— no progress in N s → abort the leg safely."),
], size=14, gap=12, mcolor=INDIGO)
footer(s)

# ---------------- OBSTACLE STRATEGY ----------------
s = slide(); header(s, "Obstacles", "Obstacle avoidance — map-based, with a sensor roadmap", color=TERRA)
sketch(s, "obstacles.png", Inches(0.55), Inches(1.85), Inches(8.4), Inches(5.0),
       caption="Figure 8 · Avoidance — implemented (map-based) + roadmap (sensors)")
glass(s, Inches(9.1), Inches(2.0), Inches(3.55), Inches(4.5), border=TERRA, ba=28)
bullets(s, Inches(9.32), Inches(2.25), Inches(3.1), Inches(4.1), [
    ("Map-based avoidance ", "— routes around configured keep-out zones (implemented, unit-tested)."),
    ("Detour planning ", "— a blocked leg gets waypoints at a safe lateral offset."),
    ("Altitude + geofence ", "— flies above ground obstacles, inside a bounded area."),
    ("Failsafe abort ", "— battery / GPS / link → RTL or LAND."),
    ("Roadmap ", "— sensor-based reactive avoidance on Copter 4.x."),
], size=13, gap=10)
footer(s)

# ---------------- 6 · STATE MACHINE ----------------
s = slide(); header(s, "Mission lifecycle", "A 12-state flight machine — every step logged")
sketch(s, "states.png", Inches(0.55), Inches(1.85), Inches(8.5), Inches(5.0),
       caption="Figure 3 · Mission state machine")
glass(s, Inches(9.2), Inches(2.0), Inches(3.45), Inches(4.4), border=OLIVE, ba=28)
bullets(s, Inches(9.45), Inches(2.25), Inches(2.95), Inches(4.0), [
    ("Happy path ", "— arm → take off → enroute → hover → RTL → land → done."),
    ("Any phase ", "can divert safely to ABORTED or FAILED."),
    ("Abort blocks ", "until disarm — never hand off an airborne drone."),
], size=14, gap=13, mcolor=OLIVE)
footer(s, 6)

# ---------------- 7 · SEQUENCE ----------------
s = slide(); header(s, "Dispatch flow", "From network trigger to touchdown")
sketch(s, "seq.png", Inches(0.55), Inches(1.85), Inches(8.0), Inches(5.0),
       caption="Figure 4 · Mission message sequence")
glass(s, Inches(8.7), Inches(2.0), Inches(3.95), Inches(4.4), border=INDIGO, ba=28)
bullets(s, Inches(8.95), Inches(2.25), Inches(3.5), Inches(4.0), [
    ("1 · Trigger ", "— validated at the edge vs the geofence."),
    ("2 · Queue ", "— bounded, strictly serial."),
    ("3 · Execute ", "— connect, GPS lock, arm, climb."),
    ("4 · Verify ", "— each mode change confirmed."),
    ("5 · Return ", "— hover, RTL, land, log incident."),
], size=14, gap=12, mcolor=INDIGO)
footer(s, 7)

# ---------------- 8 · PATENT 1 ----------------
s = slide(); header(s, "Intellectual property · Patent 1",
                    "Verified flight-mode delivery + landing-interlocked dispatch", color=MUST)
text(s, Inches(0.74), Inches(1.9), Inches(11.85), Inches(0.55),
     "Treat the autopilot as untrusted; prove every command; never overlap missions.",
     size=18, color=SOFT, italic=True)
bullets(s, Inches(0.74), Inches(2.6), Inches(11.85), Inches(3.4), [
    ("Confirmed setter ", "— issue via high-level API, then re-issue as COMMAND_LONG / DO_SET_MODE and SET_MODE until telemetry confirms."),
    ("Cross-action fallback ", "— if RTL won't confirm, command LAND through the same verified routine, and vice versa."),
    ("Abort guarantee ", "— every abnormal exit blocks until the vehicle reports disarmed (or a 240 s bound)."),
    ("Pre-flight interlock ", "— refuse to arm/take off while armed — violation becomes impossible, not just improbable."),
    ("Edge geofence ", "— reject out-of-fence targets & waypoints before flight."),
], size=16.5, gap=13, mcolor=MUST)
glass(s, Inches(0.74), Inches(6.0), Inches(11.85), Inches(0.85), border=MUST, ba=26)
text(s, Inches(1.05), Inches(6.12), Inches(11.3), Inches(0.6),
     "8 claims · IPC B64U 10/13, G05D 1/00, G08G 5/00 · IPO Form 2 complete-specification draft",
     size=14, color=MUST, font=F_BODY, italic=True, anchor=MSO_ANCHOR.MIDDLE)
footer(s, 8)

# ---------------- 9 · VERIFIED SETTER (code on paper) ----------------
s = slide(); header(s, "Patent 1 · mechanism", "The verified mode setter, in practice")
glass(s, Inches(0.74), Inches(1.95), Inches(7.5), Inches(4.6), fa=72, ba=24, border=INK)
code = [
 ("# issue, then confirm from the autopilot's OWN telemetry", {"color": MUTED}),
 ("v.mode = target              # 1) high-level setter", {"color": INK}),
 ("while now < deadline:", {"color": INDIGO}),
 ("    if v.mode.name == target:   # confirmed!", {"color": OLIVE}),
 ("        return True", {"color": OLIVE}),
 ("    raw_set_mode(target)     # 2) COMMAND_LONG", {"color": INDIGO}),
 ("    set_mode_send(target)    # 3) SET_MODE", {"color": INDIGO}),
 ("# emergency? if RTL fails, fall back to LAND", {"color": MUTED}),
 ("if not confirmed(action):", {"color": INDIGO}),
 ("    set_mode_confirmed(other)   # cross-action", {"color": TERRA}),
]
text(s, Inches(1.05), Inches(2.25), Inches(7.0), Inches(4.1), code,
     size=14, font="Consolas", spacing=1.32)
glass(s, Inches(8.5), Inches(1.95), Inches(4.1), Inches(4.6), border=INK, ba=28)
bullets(s, Inches(8.75), Inches(2.2), Inches(3.6), Inches(4.2), [
    ("Idempotent ", "— redundant re-issue is harmless."),
    ("Zero client trust ", "— only telemetry confirms."),
    ("Emergencies too ", "— abort path as robust as nominal."),
    ("Proven in SITL ", "— silent-mode-reject bug mitigated."),
], size=14.5, gap=13)
footer(s, 9)

# ---------------- 10 · PATENT 2 ----------------
s = slide(); header(s, "Intellectual property · Patent 2",
                    "Debounced, severity-ordered failsafe arbitration", color=MUST)
text(s, Inches(0.74), Inches(1.9), Inches(11.85), Inches(0.55),
     "Many hazards, evaluated at 1 Hz, resolve to exactly one calm recovery action.",
     size=18, color=SOFT, italic=True)
bullets(s, Inches(0.74), Inches(2.6), Inches(11.85), Inches(3.5), [
    ("Monotone severity ", "— actions ordered NONE < RTL < LAND; a LAND demand is never downgraded."),
    ("Debounced GPS ", "— fires only after N consecutive bad samples (reset on recovery); a glitch never lands the drone."),
    ("Mid-recovery escalation ", "— a critical-battery LAND overrides an in-progress RTL."),
    ("Fire-once discipline ", "— each hazard emits one event per mission (re-fire only to escalate)."),
    ("Above firmware ", "— layered on unmodified ArduPilot/PX4 failsafes."),
], size=16.5, gap=13, mcolor=MUST)
glass(s, Inches(0.74), Inches(6.0), Inches(11.85), Inches(0.85), border=MUST, ba=26)
text(s, Inches(1.05), Inches(6.12), Inches(11.3), Inches(0.6),
     "8 claims · pinned by a 37-case unit suite · IPC G05D 1/00, B64D 45/00, B64U 2201/10",
     size=14, color=MUST, font=F_BODY, italic=True, anchor=MSO_ANCHOR.MIDDLE)
footer(s, 10)

# ---------------- 11 · FAILSAFE TREE ----------------
s = slide(); header(s, "Patent 2 · mechanism", "The hazard tree: condition → threshold → action")
sketch(s, "failsafe.png", Inches(0.6), Inches(1.85), Inches(8.4), Inches(5.0),
       caption="Figure 5 · Failsafe decision tree")
glass(s, Inches(9.15), Inches(2.0), Inches(3.5), Inches(4.4), border=OLIVE, ba=28)
bullets(s, Inches(9.4), Inches(2.25), Inches(3.05), Inches(4.0), [
    ("RTL ", "— battery low, geofence, timeout, link loss."),
    ("LAND ", "— battery critical, GPS lost (can't navigate)."),
    ("Debounce ", "— GPS needs N bad samples to fire."),
], size=14, gap=13, mcolor=OLIVE)
footer(s, 11)

# ---------------- 12 · INTERLOCK ----------------
s = slide(); header(s, "Safety model", "Five gates from network trigger to landing")
sketch(s, "interlock.png", Inches(0.7), Inches(2.0), Inches(11.9), Inches(2.9))
glass(s, Inches(0.74), Inches(5.2), Inches(11.85), Inches(1.0), border=OLIVE, ba=28)
rect(s, Inches(0.74), Inches(5.2), Inches(0.12), Inches(1.0), fill=OLIVE, rounded=False)
text(s, Inches(1.1), Inches(5.4), Inches(11.3), Inches(0.6),
     [("The guarantee:  ", {"bold": True, "color": INK, "font": F_TITLE, "size": 17}),
      ("even if four gates fail, no mission ever starts against an airborne drone.",
       {"color": SOFT, "size": 17, "italic": True})], anchor=MSO_ANCHOR.MIDDLE)
text(s, Inches(0.7), Inches(4.85), Inches(11.9), Inches(0.35),
     "Figure 6 · The five-gate safety interlock", size=13, color=MUTED,
     font=F_BODY, italic=True, align=PP_ALIGN.CENTER)
footer(s, 12)

# ---------------- 13 · VALIDATION ----------------
s = slide(); header(s, "Validation", "Proven in software-in-the-loop", color=OLIVE)
sketch(s, "traj.png", Inches(0.7), Inches(1.9), Inches(6.9), Inches(4.9),
       caption="Figure 7 · Reconstructed flight path")
kp = [("8 / 8", "acceptance checks pass", OLIVE), ("0.4 m", "terminal accuracy (5 m tol.)", TERRA),
      ("37", "safety unit tests, ms-fast", INDIGO), ("328 s", "full autonomous mission", MUST)]
y = Inches(2.0)
for v, l, c in kp:
    glass(s, Inches(8.0), y, Inches(4.6), Inches(1.05), border=c, ba=26)
    rect(s, Inches(8.0), y, Inches(0.12), Inches(1.05), fill=c, rounded=False)
    text(s, Inches(8.35), Emu(int(y) + int(Inches(0.1))), Inches(2.2), Inches(0.85),
         v, size=31, color=c, font=F_TITLE, bold=True)
    text(s, Inches(10.4), y, Inches(2.05), Inches(1.05), l, size=14, color=SOFT,
         font=F_BODY, italic=True, anchor=MSO_ANCHOR.MIDDLE, spacing=1.05)
    y = Emu(int(y) + int(Inches(1.18)))
footer(s, 13)

# ---------------- 14 · TECH STACK ----------------
s = slide(); header(s, "Engineering", "Technology stack")
groups = [
    ("Flight core", TERRA, ["Python 3.11", "DroneKit + pymavlink", "MAVLink protocol", "ArduPilot SITL"]),
    ("Service", OLIVE, ["FastAPI + Pydantic v2", "Uvicorn ASGI", "WebSocket telemetry", "SQLite store"]),
    ("Dashboard", INDIGO, ["React 18 + Vite 5", "react-leaflet map", "Live path + log", "Follow-drone toggle"]),
    ("Quality / ops", MUST, ["37 pytest tests", "E2E SITL acceptance", "Docker compose", "API-key + geofence"]),
]
x = Inches(0.74); w = Inches(2.83)
for title, c, items in groups:
    glass(s, x, Inches(2.0), w, Inches(4.3), border=c, ba=28)
    rect(s, Emu(int(x) + int(Inches(0.22))), Inches(2.22),
         Emu(int(w) - int(Inches(0.44))), Pt(2.2), fill=c)
    text(s, Emu(int(x) + int(Inches(0.26))), Inches(2.42), Emu(int(w) - int(Inches(0.5))),
         Inches(0.5), title, size=17, color=c, font=F_TITLE, bold=True)
    bullets(s, Emu(int(x) + int(Inches(0.26))), Inches(3.05), Emu(int(w) - int(Inches(0.46))),
            Inches(3.0), items, size=14.5, gap=13, mcolor=c)
    x = Emu(int(x) + int(Inches(3.0)))
footer(s, 14)

# ---------------- 15 · ROADMAP ----------------
s = slide(); header(s, "Roadmap", "Next steps & Claude Code integrations")
glass(s, Inches(0.74), Inches(2.0), Inches(5.7), Inches(4.4), border=TERRA, ba=28)
rect(s, Emu(int(Inches(0.96))), Inches(2.22), Inches(5.26), Pt(2.2), fill=TERRA)
text(s, Inches(1.0), Inches(2.42), Inches(5.3), Inches(0.5), "Product hardening",
     size=17, color=TERRA, font=F_TITLE, bold=True)
bullets(s, Inches(1.0), Inches(3.05), Inches(5.3), Inches(3.2), [
    ("Migrate off DroneKit ", "→ MAVSDK / ArduPilot 4.x."),
    ("Failsafe-injection tests ", "— scripted SITL cuts."),
    ("Multi-leg routes ", "— full waypoint list up front."),
    ("RTK GPS ", "for ~0.1 m real-world accuracy."),
], size=15, gap=13)
glass(s, Inches(6.64), Inches(2.0), Inches(5.98), Inches(4.4), border=INDIGO, ba=28)
rect(s, Emu(int(Inches(6.86))), Inches(2.22), Inches(5.54), Pt(2.2), fill=INDIGO)
text(s, Inches(6.9), Inches(2.42), Inches(5.5), Inches(0.5), "Claude Code integrations",
     size=17, color=INDIGO, font=F_TITLE, bold=True)
bullets(s, Inches(6.9), Inches(3.05), Inches(5.5), Inches(3.2), [
    ("/fly slash command ", "— trigger → telemetry → cancel."),
    ("Save-hook ", "— auto-run unit tests on edits."),
    ("GitHub Actions CI ", "— run the safety suite per push."),
    ("MCP telemetry server ", "— query live state in chat."),
], size=15, gap=13, mcolor=INDIGO)
footer(s, 15)

# ---------------- 16 · IP STATUS ----------------
s = slide(); header(s, "Filing", "Copyright vs. patent — what protects what", color=MUST)
glass(s, Inches(0.74), Inches(2.0), Inches(5.7), Inches(2.5), border=OLIVE, ba=28)
rect(s, Emu(int(Inches(0.96))), Inches(2.22), Inches(5.26), Pt(2.2), fill=OLIVE)
text(s, Inches(1.0), Inches(2.42), Inches(5.3), Inches(0.5), "Copyright",
     size=18, color=OLIVE, font=F_TITLE, bold=True)
text(s, Inches(1.0), Inches(3.0), Inches(5.3), Inches(1.4),
     "Protects the expression — source code, docs, diagrams, the paper. Automatic on "
     "authorship; registration is cheap, fast, and yours to file now.",
     size=15, color=SOFT, spacing=1.16)
glass(s, Inches(6.64), Inches(2.0), Inches(5.98), Inches(2.5), border=MUST, ba=28)
rect(s, Emu(int(Inches(6.86))), Inches(2.22), Inches(5.54), Pt(2.2), fill=MUST)
text(s, Inches(6.9), Inches(2.42), Inches(5.5), Inches(0.5), "Patent",
     size=18, color=MUST, font=F_TITLE, bold=True)
text(s, Inches(6.9), Inches(3.0), Inches(5.5), Inches(1.4),
     "Protects the inventions — the verified setter, landing interlock, failsafe arbiter. "
     "Two IPO Form-2 drafts exist; needs a novelty search + a registered agent.",
     size=15, color=SOFT, spacing=1.16)
glass(s, Inches(0.74), Inches(4.7), Inches(11.88), Inches(1.75), border=MUST, ba=26)
text(s, Inches(1.0), Inches(4.88), Inches(11.3), Inches(0.4), "Recommended sequence",
     size=15, color=MUST, font=F_TITLE, italic=True)
bullets(s, Inches(1.0), Inches(5.32), Inches(11.3), Inches(1.0), [
    ("Now ", "— register copyright on the codebase + paper (immediate, low-cost)."),
    ("Then ", "— run IPO InPASS / WIPO novelty search against both patent drafts."),
    ("Then ", "— file provisional patents via a registered agent before public disclosure."),
], size=15, gap=7, mcolor=MUST)
footer(s, 16)

# ---------------- 17 · CLOSING ----------------
s = slide()
glass(s, Inches(1.4), Inches(1.85), Inches(10.5), Inches(3.8), fa=58, ba=22)
text(s, Inches(1.9), Inches(2.3), Inches(9.5), Inches(0.5),
     "Autonomous · Verified · Safe", size=16, color=TERRA, font=F_TITLE, italic=True)
text(s, Inches(1.85), Inches(2.8), Inches(9.7), Inches(1.2),
     "Drone Safety System", size=50, color=INK, font=F_TITLE, bold=True)
rect(s, Inches(1.92), Inches(3.95), Inches(2.7), Pt(2.4), fill=TERRA)
text(s, Inches(1.9), Inches(4.15), Inches(9.5), Inches(0.6),
     "A provable safety layer for autonomous UAV dispatch.", size=21, color=SOFT, italic=True)
text(s, Inches(1.9), Inches(4.95), Inches(9.5), Inches(0.4),
     "github.com/SV-1411/drone", size=17, color=INDIGO, font=F_BODY, italic=True)
text(s, Inches(1.4), Inches(6.0), Inches(10.5), Inches(0.4),
     "Built with a companion-computer architecture · DroneKit · MAVLink · ArduPilot",
     size=13, color=MUTED, font=F_BODY, italic=True, align=PP_ALIGN.CENTER)

prs.save(os.path.join(HERE, "DroneSafetySystem.pptx"))
print("saved", os.path.join(HERE, "DroneSafetySystem.pptx"), "slides:", len(prs.slides._sldIdLst))
